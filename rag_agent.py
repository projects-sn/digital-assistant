# rag_agent.py
import os
import re
import json
import hashlib
from datetime import datetime
from typing import List, Tuple, Optional, Any

import pandas as pd
import openai

from langchain_openai import OpenAIEmbeddings
from langchain_text_splitters import TokenTextSplitter
from langchain_core.documents import Document
from langchain_community.retrievers import BM25Retriever
from langchain_community.vectorstores import DocArrayInMemorySearch
from langchain_openai import ChatOpenAI
from langchain_core.prompts import PromptTemplate
from langchain_core.retrievers import BaseRetriever

try:
    from docx import Document as DocxDocument  # type: ignore[import]
except ImportError:  # pragma: no cover
    DocxDocument = None  # type: ignore[assignment]

try:
    from pptx import Presentation  # type: ignore[import]
except ImportError:  # pragma: no cover
    Presentation = None  # type: ignore[assignment]

from dotenv import load_dotenv
load_dotenv()

# === Глобальные структуры (обновляются при переиндексации) ===
_split_docs: List[Document] = []
_embeddings = None
_vectorstore = None
_bm25: Optional[BM25Retriever] = None


# ======================
# Вспомогательные функции
# ======================

def _safe_meta_date(s: Optional[str]) -> str:
    if not s:
        return "неизвестно"
    try:
        if re.match(r"\d{2}\.\d{2}\.\d{4}$", str(s)):
            return str(s)
        dt = datetime.fromisoformat(str(s))
        return dt.strftime("%d.%m.%Y")
    except Exception:
        return str(s)

def _mk_doc_id(text: str, path: Optional[str] = None) -> str:
    base = (path or "") + "|" + text[:200]
    return hashlib.md5(base.encode("utf-8")).hexdigest()[:10]

def _as_documents_from_text(text: str, *, source: str, date: Optional[str] = None, label: Optional[str] = None) -> List[Document]:
    meta = {
        "source": source,
        "date": _safe_meta_date(date),
        "label": label or os.path.basename(source),
        "doc_id": _mk_doc_id(text, source),
    }
    page = f"Источник: {meta['label']}\nДата: {meta['date']}\nФайл: {source}\n\n{text}"
    return [Document(page_content=page, metadata=meta)]


# ======================
# Чтение разных форматов
# ======================

def _read_txt(path: str) -> str:
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()

def _read_md(path: str) -> str:
    return _read_txt(path)

def _read_pdf(path: str) -> str:
    from pypdf import PdfReader
    reader = PdfReader(path)
    parts = []
    for page in reader.pages:
        try:
            parts.append(page.extract_text() or "")
        except Exception:
            parts.append("")
    return "\n\n".join(parts).strip()

def _read_docx(path: str) -> str:
    if DocxDocument is None:
        raise ImportError("python-docx не установлен")
    d = DocxDocument(path)
    return "\n".join(p.text for p in d.paragraphs)

def _read_pptx(path: str) -> str:
    if Presentation is None:
        raise ImportError("python-pptx не установлен")
    prs = Presentation(path)
    chunks = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                chunks.append(shape.text)
    return "\n\n".join(chunks).strip()


# ======================
# Специальные лоадеры
# ======================

def _load_from_csv(csv_path: str) -> List[Document]:
    """Старый источник: meeting_summaries.csv (колонки: Дата, Документ, Сводка)."""
    if not os.path.exists(csv_path):
        return []
    df = pd.read_csv(csv_path).dropna(subset=["Сводка"])
    docs: List[Document] = []
    for _, row in df.iterrows():
        date = _safe_meta_date(str(row.get("Дата", "")))
        docno = str(row.get("Документ", ""))
        text = str(row.get("Сводка", ""))

        meta = {
            "doc_id": str(docno),
            "date": date,
            "source": os.path.basename(csv_path),
            "label": f"Документ №{docno}",
        }
        page = f"Дата: {date}\nДокумент №: {docno}\n\n{text}"
        docs.append(Document(page_content=page, metadata=meta))
    return docs

def _load_from_presentations_json(path: str) -> List[Document]:
    """
    Ожидаем структуру:
    {
      "presentations": [
        {"filename": "...", "pages": [{"page": 1, "text": "..."}, ...]}
      ]
    }
    Каждую страницу превращаем в Document с аккуратными метаданными.
    """
    if not os.path.exists(path):
        return []
    docs: List[Document] = []
    with open(path, "r", encoding="utf-8") as f:
        data = json.load(f)

    for pres in (data.get("presentations") or []):
        fname = str(pres.get("filename") or "presentation")
        for p in (pres.get("pages") or []):
            page_no = p.get("page")
            text = (p.get("text") or "").strip()
            if not text:
                continue
            meta = {
                "source": path,
                "label": f"{fname} — стр. {page_no}",
                "date": "неизвестно",
                "doc_id": _mk_doc_id(f"{fname}:{page_no}", path),
            }
            page = f"Источник: {fname}\nСтраница: {page_no}\nФайл: {path}\n\n{text}"
            docs.append(Document(page_content=page, metadata=meta))
    return docs

def _load_from_folder(folder: str = "data") -> List[Document]:
    """Сканируем папку и читаем txt/md/pdf/docx/pptx."""
    docs: List[Document] = []
    if not os.path.isdir(folder):
        return docs
    for root, _, files in os.walk(folder):
        for name in files:
            path = os.path.join(root, name)
            ext = os.path.splitext(name)[1].lower()
            try:
                if ext == ".txt":
                    text = _read_txt(path)
                elif ext in (".md", ".markdown"):
                    text = _read_md(path)
                elif ext == ".pdf":
                    text = _read_pdf(path)
                elif ext == ".docx":
                    text = _read_docx(path)
                elif ext == ".pptx":
                    text = _read_pptx(path)
                else:
                    continue  # пропускаем неподдерживаемые типы

                if text.strip():
                    docs.extend(_as_documents_from_text(text, source=path, date=None, label=None))
            except Exception as e:
                print(f"[WARN] Не удалось прочитать {path}: {e}")
    return docs


# ======================
# Индексация корпуса
# ======================

def _build_indexes(docs: List[Document]):
    global _split_docs, _embeddings, _vectorstore, _bm25
    if not docs:
        _split_docs = []
        _vectorstore = None
        _bm25 = None
        return

    splitter = TokenTextSplitter(chunk_size=400, chunk_overlap=40)
    _split_docs = splitter.split_documents(docs)

    _embeddings = OpenAIEmbeddings(
        model="text-embedding-3-small",
        openai_api_key=os.getenv("OPENAI_API_KEY")
    )
    _vectorstore = DocArrayInMemorySearch.from_documents(_split_docs, _embeddings)

    _bm25 = BM25Retriever.from_documents(_split_docs)
    _bm25.k = 4

def reindex_corpus() -> int:
    """
    Полная переиндексация: CSV + JSON презентаций + папка data/.
    """
    base_docs: List[Document] = []
    # 1) CSV (совместимость с текущей схемой)
    base_docs += _load_from_csv("meeting_summaries.csv")
    # 2) JSON с текстами презентаций (если положен в uploads)
    base_docs += _load_from_presentations_json(os.path.join("data", "uploads", "all_presentations.json"))
    # 3) Всё остальное из папки (txt/md/pdf/docx/pptx, включая ваши три стенограммы .txt)
    base_docs += _load_from_folder("data")

    _build_indexes(base_docs)
    return len(base_docs)

def add_files_and_reindex(saved_paths: List[str]) -> int:
    """
    Вызывается после загрузки файлов через UI. Просто запускаем полную переиндексацию.
    """
    return reindex_corpus()


# ======================
# Поисковая логика (гибрид)
# ======================

def _call_retriever(retriever: BM25Retriever, query: str) -> List[Document]:
    """
    Совместимость с разными версиями langchain-community:
    часть реализаций предоставляет get_relevant_documents, часть — invoke().
    """
    if hasattr(retriever, "get_relevant_documents"):
        docs = retriever.get_relevant_documents(query)
    else:
        docs = retriever.invoke(query)  # type: ignore[arg-type]
    if not isinstance(docs, list):
        return []
    return docs

def _hybrid_retrieve(query: str, k: int = 4) -> List[Document]:
    if not _split_docs or _vectorstore is None or _bm25 is None:
        reindex_corpus()
    vector_docs = _vectorstore.similarity_search(query, k=k)
    bm25_docs = _call_retriever(_bm25, query)
    all_docs = vector_docs + bm25_docs
    unique = {}
    for doc in all_docs:
        key = doc.page_content[:200]
        if key not in unique:
            unique[key] = doc
    return list(unique.values())[:k]

def extract_date_and_doc_id(query: str) -> Tuple[Optional[str], Optional[str]]:
    date_match = re.search(r"(\d{2}\.\d{2}\.\d{4})", query)
    doc_match = re.search(r"(документ|встреча)\s*№?\s*(\d+)", query.lower())
    date = date_match.group(1) if date_match else None
    doc_id = doc_match.group(2) if doc_match else None
    return date, doc_id

class CustomRetriever(BaseRetriever):
    def _retrieve(self, query: str) -> List[Document]:
        if not _split_docs:
            reindex_corpus()

        date, doc_id = extract_date_and_doc_id(query)

        if date or doc_id:
            filtered = [
                d for d in _split_docs
                if (not date or d.metadata.get("date") == date)
                and (not doc_id or d.metadata.get("doc_id") == doc_id)
            ]
            if filtered:
                temp_vs = DocArrayInMemorySearch.from_documents(filtered, _embeddings)
                temp_bm = BM25Retriever.from_documents(filtered); temp_bm.k = 4
                vec = temp_vs.similarity_search(query, k=4)
                kw = _call_retriever(temp_bm, query)
                uniq = {}
                for dd in vec + kw:
                    key = dd.page_content[:200]
                    if key not in uniq:
                        uniq[key] = dd
                return list(uniq.values())[:4]

        return _hybrid_retrieve(query, k=4)

    # BaseRetriever API
    def _get_relevant_documents(self, query: str, *, run_manager=None) -> List[Document]:
        return self._retrieve(query)

    def get_relevant_documents(self, query: str) -> List[Document]:
        return self._retrieve(query)


# ======================
# RAG-цепочка (кастомная реализация без langchain.chains)
# ======================

class SimpleRetrievalQA:
    def __init__(self, llm: ChatOpenAI, *, retriever: BaseRetriever, prompt: PromptTemplate):
        self.llm = llm
        self.retriever = retriever
        self.prompt = prompt

    def _run(self, question: str) -> dict[str, Any]:
        docs = self.retriever.get_relevant_documents(question)
        context = "\n\n".join(doc.page_content for doc in docs)
        prompt_str = self.prompt.format(context=context, question=question)
        response = self.llm.invoke(prompt_str)
        text = getattr(response, "content", str(response))
        return {"result": text.strip(), "source_documents": docs}

    def invoke(self, inputs: dict[str, Any]) -> dict[str, Any]:
        question = inputs.get("query") or inputs.get("question")
        if not question:
            raise ValueError("SimpleRetrievalQA ожидает ключ 'query' или 'question'")
        return self._run(question)

    def __call__(self, inputs: dict[str, Any]) -> dict[str, Any]:
        return self.invoke(inputs)


QA_PROMPT = PromptTemplate(
    input_variables=["context", "question"],
    template="""
Ты — помощник, анализирующий материалы встреч, стенограммы и презентации.
Отвечай по делу, опираясь на контекст. В конце кратко перечисли источники (название/файл и дату).

Контекст:
{context}

Вопрос:
{question}

Ответ:
"""
)

def build_rag_chain():
    # Ленивая индексация
    if not _split_docs:
        reindex_corpus()

    llm = ChatOpenAI(
        model_name="gpt-4",
        temperature=0.2,
        max_tokens=2048,
        openai_api_key=os.getenv("OPENAI_API_KEY"),
    )
    retriever = CustomRetriever()
    return SimpleRetrievalQA(llm, retriever=retriever, prompt=QA_PROMPT)
